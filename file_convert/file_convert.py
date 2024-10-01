import json
import os
from pathlib import Path
from urllib.parse import urljoin
import requests
import fitz
import pandas as pd
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
import re

from common.local_models import query_for_local_model
from common.tools import generate_hash_value, extract_keywords, save_as_json, save_as_md
from filter.config import Filter_Config
from config import Convert_Config

FilterConfig = Filter_Config()
ConvertConfig = Convert_Config()
OUTPUT_JSON_PATH:str=FilterConfig.PURIED_JSON_PATH+FilterConfig.SCHOOL_SIMPLE
OUTPUT_FMD_PATH:str=ConvertConfig.COMMON_OUTPUT_PATH
class FileConvert:
    """封装所有文件读取操作。"""
    md_image_pattern = re.compile(r'\[([^\]]+)\]\(([a-zA-Z0-9:/._~#-]+)?\)')
    html_image_pattern = re.compile(r'<img\s+[^>]*?src=["\']([^"\']*)["\'][^>]*>')

    def __init__(self):
        self.md_suffix = '.md'
        self.text_suffix = ['.txt', '.text']
        self.excel_suffix = ['.xlsx', '.xls', '.csv']
        self.pdf_suffix = '.pdf'
        self.ppt_suffix = '.pptx'
        self.html_suffix = ['.html', '.htm', '.shtml', '.xhtml']
        self.word_suffix = ['.docx', '.doc']
        self.normal_suffix = [self.md_suffix
                              ] + self.text_suffix + self.excel_suffix + [
                                 self.pdf_suffix
                             ] + self.word_suffix + [self.ppt_suffix
                                                     ] + self.html_suffix

    def convert(self, filepath: str,file_name:str):
        """转换文件格式。"""
        if not os.path.isabs(filepath):
            raise ValueError("File path must be an absolute path.")

        file_type = self.get_type(filepath)
        file_name = os.path.splitext(file_name)[0]
        if file_type == "md":
            return filepath
        elif file_type == "pdf":
            return self.convert_pdf(filepath,file_name)
        elif file_type == "excel":
            return self.convert_excel(filepath,file_name)
        elif file_type == "word":
            return self.convert_word(filepath,file_name)
        elif file_type == "ppt":
            return self.convert_ppt(filepath,file_name)
        elif file_type == "html":
            return self.convert_html(filepath,file_name)
        elif file_type == "text":
            return self.convert_text(filepath,file_name)
        else:
            return None

    def get_type(self, filepath: str):
        """根据URI后缀获取文件类型。"""
        filepath = filepath.lower()
        if filepath.endswith(self.pdf_suffix):
            return 'pdf'

        if filepath.endswith(self.md_suffix):
            return 'md'

        if filepath.endswith(self.ppt_suffix):
            return 'ppt'

        for suffix in self.text_suffix:
            if filepath.endswith(suffix):
                return 'text'

        for suffix in self.word_suffix:
            if filepath.endswith(suffix):
                return 'word'

        for suffix in self.excel_suffix:
            if filepath.endswith(suffix):
                return 'excel'

        for suffix in self.html_suffix:
            if filepath.endswith(suffix):
                return 'html'

        return None

    def save_image(self, image, image_name,img_path):
        """保存图片到images目录。"""
        images_dir = Path(img_path)
        images_dir.mkdir(parents=True, exist_ok=True)
        image_path = images_dir / image_name

        with open(image_path, 'wb') as img_file:
            img_file.write(image)

    def convert_pdf(self, filepath: str,pdf_name):
        """读取PDF文件并转换为Markdown格式。"""
        text = ''
        output_path = os.path.join(OUTPUT_FMD_PATH,"pdf2md",ConvertConfig.SCHOOL_SIMPLE)
        img_path = os.path.join(output_path,pdf_name,"images")
        with fitz.open(filepath) as pages:
            for page_num, page in enumerate(pages, start=1):
                text += page.get_text()
                for img_index, img in enumerate(page.get_images(full=True), start=1):
                    xref = img[0]
                    base_image = pages.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_name = f"page_{page_num}_img_{img_index}.png"
                    self.save_image(image_bytes, image_name,img_path)
                    text += f"![image](images/{image_name})\n\n"
        save_as_md(output_path,pdf_name,text)
        return text

    def convert_excel(self, filepath: str,excel_name:str):
        """读取Excel文件并转换为Markdown格式。"""
        table = None
        output_path = os.path.join(OUTPUT_FMD_PATH, "xlsx2md",ConvertConfig.SCHOOL_SIMPLE)
        if filepath.endswith('.csv'):
            table = pd.read_csv(filepath)
        else:
            table = pd.read_excel(filepath)
        if table is None:
            return ''
        md_content = table.to_markdown(index=False)
        save_as_md(output_path, excel_name, md_content)

    def convert_word(self, filepath: str,word_name:str) -> str:
        """读取Word文档并转换为Markdown格式。"""
        doc = Document(filepath)
        markdown_content = ""
        output_path = os.path.join(OUTPUT_FMD_PATH,"docx2md",ConvertConfig.SCHOOL_SIMPLE)
        img_path = os.path.join(output_path,word_name,"images")
        for paragraph in doc.paragraphs:
            markdown_content += f"{paragraph.text}\n\n"
        for i, shape in enumerate(doc.inline_shapes, start=1):
            if shape.type == 3:  # InlineShapePicture
                image = shape._inline.graphic.graphicData.pic.blipFill.blip.embed
                image_bytes = doc.part.related_parts[image].blob
                image_name = f"image_{i}.png"
                self.save_image(image_bytes, image_name,img_path)
                markdown_content += f"![image](images/{image_name})\n\n"
        save_as_md(output_path, word_name, markdown_content)
        return markdown_content

    def convert_ppt(self, filepath: str,ppt_name:str) -> str:
        """读取PPT文件并转换为Markdown格式。"""
        presentation = Presentation(filepath)
        markdown_content = ""
        output_path = os.path.join(OUTPUT_FMD_PATH,"pptx2md",ConvertConfig.SCHOOL_SIMPLE)
        img_path = os.path.join(output_path,ppt_name,"images")
        for slide_num, slide in enumerate(presentation.slides, start=1):
            slide_title = slide.shapes.title.text if slide.shapes.title else "Slide Title"
            markdown_content += f"# {slide_title}\n\n"
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    markdown_content += f"{paragraph.text}\n\n"
                if shape.shape_type == 13:  # Picture
                    image = shape.image
                    image_bytes = image.blob
                    image_name = f"slide_{slide_num}_img.png"
                    self.save_image(image_bytes, image_name,img_path)
                    markdown_content += f"![image](images/{image_name})\n\n"
        save_as_md(output_path, ppt_name, markdown_content)
        return markdown_content

    # HTML清洗，转换及保存
    def convert_html(self, filepath: str,html_name:str) -> str:
        """读取HTML文件并转换为Markdown格式。"""
        with open(filepath, 'r', encoding='utf-8') as file:
            soup = BeautifulSoup(file, 'html.parser')
        markdown_content = soup.get_text()
        # 正则表达式匹配多个换行符，用一个换行符全替换
        pat = r'\n{2,}'
        markdown_content = re.sub(pat, '\n', markdown_content)
        attrs = self.attr_process(markdown_content,html_name)
        img_path = os.path.join(OUTPUT_JSON_PATH,attrs['title']+"/images")
        for i, img in enumerate(soup.find_all('img'), start=1):
            # 补充URL恢复逻辑
            img_url = img['src']
            html_link = self.get_html_link(html_name)
            img_url = self.img_url_recover(img_url,html_link)
            image_bytes = requests.get(img_url).content
            image_name = f"image_{i}.png"
            self.save_image(image_bytes, image_name,img_path)
            markdown_content += f"![image](images/{image_name})\n\n"

        hash_value = generate_hash_value(markdown_content)
        keywords = extract_keywords(markdown_content,self.md_image_pattern)
        # 保存到每个md子节目，md与json是平级的，img都在下一级
        save_as_json(attrs["title"], attrs["publish_date"], keywords, attrs["category"], markdown_content, hash_value)
        save_as_md(OUTPUT_JSON_PATH,attrs["title"], markdown_content)
        return markdown_content

    def convert_text(self, filepath: str,txt_name:str) -> str:
        output_path = os.path.join(OUTPUT_FMD_PATH, "txt2md", ConvertConfig.SCHOOL_SIMPLE)
        """读取文本文件并转换为Markdown格式。"""
        with open(filepath, 'r', encoding='utf-8') as file:
            markdown_content = file.read()
        save_as_md(output_path, txt_name, markdown_content)

    def get_html_link(self, html_name) -> str:
        with open("url_mapping.json", 'r') as f:
            url_table = json.load(f)
        return url_table[html_name]
    def join_relative_url(self,base_url: str, relative_url: str) -> str:
        return urljoin(base_url, relative_url)
    def img_url_recover(self,ab_url:str,html_link:str)->str:
        if ab_url.startswith('http') or ab_url.startswith('https'):
            return ab_url
        else:
            complete_url = self.join_relative_url(html_link, ab_url)
            return complete_url
    def generate_attrs(self,content) -> dict:
        data = {
            "model": FilterConfig.LOCAL_MODEL,
            "prompt":
                FilterConfig.PROMPT_OF_ATTRS + content,
            "stream": False
        }
        attrs = query_for_local_model(data)
        response = json.loads(attrs.text)['response']
        attrs = json.loads(response)  # TODO bug
        return attrs
    def attr_process(self,content: str, html_name) -> dict[str:str]:
        """
            使用markdownify库(不借助大模型效率更高)，将html转为md，过程中解析表格
            接着，我们开始按照指定的格式为每一篇文章，构建相应的JSON
            最后，使用上一步中本地大模型生成的标题对其命名后，保存为json最终文件。
            为了日志记录，我们需要返回文章标题
        """
        attrs = self.generate_attrs(content)

        return attrs